from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Color palette ──────────────────────────────────────────────
NAVY    = RGBColor(0x1E, 0x3A, 0x5F)
BLUE    = RGBColor(0x25, 0x63, 0xEB)
LBLUE   = RGBColor(0xDB, 0xEA, 0xFE)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xF4, 0xF7, 0xFB)
DKGREY  = RGBColor(0x60, 0x7D, 0x9A)
RED     = RGBColor(0xDC, 0x26, 0x26)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
AMBER   = RGBColor(0xB4, 0x53, 0x09)
GREEN   = RGBColor(0x05, 0x96, 0x69)
BLACK   = RGBColor(0x1E, 0x2A, 0x3A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helpers ────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color=None, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = Pt(border_pt)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color and border_pt:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def header_bar(slide, title, subtitle=None):
    """Dark navy top bar with title."""
    box(slide, 0, 0, 13.33, 1.1, NAVY)
    txt(slide, title, 0.4, 0.15, 10, 0.6, size=24, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.72, 10, 0.35, size=12, color=RGBColor(0xA8,0xC4,0xE0))

def bullet_block(slide, heading, bullets, l, t, w,
                 head_color=NAVY, bullet_color=BLACK, head_size=14, bullet_size=12):
    """Heading + bullet list block."""
    # Heading
    hbox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.35))
    hbox.word_wrap = True
    htf = hbox.text_frame
    htf.word_wrap = True
    hp = htf.paragraphs[0]
    hr = hp.add_run()
    hr.text = heading
    hr.font.size = Pt(head_size)
    hr.font.bold = True
    hr.font.color.rgb = head_color

    cur_t = t + 0.38
    for b in bullets:
        bh = 0.28
        bbox = slide.shapes.add_textbox(Inches(l+0.15), Inches(cur_t), Inches(w-0.15), Inches(bh))
        bbox.word_wrap = True
        btf = bbox.text_frame
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        br = bp.add_run()
        br.text = "  -  " + b
        br.font.size = Pt(bullet_size)
        br.font.color.rgb = bullet_color
        cur_t += bh + 0.02
    return cur_t  # returns bottom Y

def flow_box(slide, label, l, t, w=1.7, h=0.55,
             fill=LBLUE, border=BLUE, tcolor=NAVY):
    b = box(slide, l, t, w, h, fill_color=fill, border_color=border, border_pt=1.5)
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = tcolor

def arrow(slide, l, t, w=0.4, h=0.04):
    """Simple horizontal arrow line."""
    conn = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    conn.fill.solid()
    conn.fill.fore_color.rgb = BLUE
    conn.line.fill.background()


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — TITLE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)

# Accent strip
box(s1, 0, 6.8, 13.33, 0.7, BLUE)

txt(s1, "Digital Black Board", 1.0, 1.6, 11.33, 1.2,
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s1, "Learning Management System", 1.0, 2.85, 11.33, 0.6,
    size=22, bold=False, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

divider = box(s1, 3.5, 3.55, 6.33, 0.05, fill_color=BLUE)

txt(s1, "Prototype Design Document", 1.0, 3.75, 11.33, 0.5,
    size=16, color=RGBColor(0xBF,0xDB,0xFE), align=PP_ALIGN.CENTER)

txt(s1, "TEAM 18  |  KL University  |  React 19 + Vite 7", 1.0, 4.25, 11.33, 0.4,
    size=13, color=RGBColor(0xA8,0xC4,0xE0), align=PP_ALIGN.CENTER)

txt(s1, "https://digital-blackboard.netlify.app", 1.0, 6.9, 11.33, 0.4,
    size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — PROJECT OVERVIEW
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s2 = prs.slides.add_slide(BLANK)
bg(s2, GREY)
header_bar(s2, "Project Overview", "What is Digital Black Board?")

# Left column - About
box(s2, 0.35, 1.25, 5.9, 5.75, WHITE, RGBColor(0xDD,0xE6,0xF0), 1)
bullet_block(s2, "About the Project", [
    "Digital Black Board is a web-based Learning Management System (LMS)",
    "Enables institutions to centralize academic management on one platform",
    "Supports course creation, enrollment, assignments, and communication",
    "Built with React 19 and Vite 7 using modern frontend architecture",
    "Features role-based access control with four distinct user roles",
    "Includes dark and light theme support with full mobile responsiveness",
    "Data is persisted using browser localStorage for seamless session continuity",
    "Deployed publicly on Netlify with GitHub-based auto-deployment",
], 0.55, 1.35, 5.6, head_color=NAVY, bullet_color=BLACK, head_size=15)

# Right column - Tech & Features
box(s2, 6.55, 1.25, 6.45, 2.7, WHITE, RGBColor(0xDD,0xE6,0xF0), 1)
bullet_block(s2, "Technology Stack", [
    "Frontend Framework: React 19",
    "Build Tool: Vite 7",
    "Language: JavaScript (ES2023)",
    "Styling: Custom CSS with CSS Variables",
    "State Management: React Context API",
    "Persistence: Browser localStorage",
    "Deployment: Netlify (CI/CD via GitHub)",
], 6.75, 1.35, 6.1, head_color=NAVY, bullet_color=BLACK, head_size=15)

box(s2, 6.55, 4.1, 6.45, 2.9, WHITE, RGBColor(0xDD,0xE6,0xF0), 1)
bullet_block(s2, "Key Features", [
    "Four role-based dashboards (Admin, Instructor, Creator, Student)",
    "Sliding panel login and signup with password strength indicator",
    "Course creation, publishing, and enrollment workflow",
    "Assignment submission and grading system",
    "Announcements, analytics, content library, and settings",
    "Hamburger menu navigation for mobile devices",
], 6.75, 4.2, 6.1, head_color=NAVY, bullet_color=BLACK, head_size=15)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — ROLES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s3 = prs.slides.add_slide(BLANK)
bg(s3, GREY)
header_bar(s3, "User Roles and Responsibilities", "4 roles with distinct permissions and access levels")

roles = [
    ("Admin", "admin@gmail.com", "Admin@123", RED,   RGBColor(0xFE,0xF2,0xF2), RGBColor(0xFC,0xA5,0xA5), [
        "Manage all registered users (add, edit, deactivate, delete)",
        "Create, edit, publish, and remove courses",
        "Post platform-wide announcements",
        "Create, assign, and grade assignments",
        "View platform-wide analytics and statistics",
        "Configure platform settings and contact information",
    ]),
    ("Instructor", "ins@gmail.com", "ins@123", PURPLE, RGBColor(0xF5,0xF3,0xFF), RGBColor(0xC4,0xB5,0xFD), [
        "View and manage own created courses",
        "Add and manage assignments for own courses",
        "Track student submissions and performance",
        "View course-level analytics",
        "Read platform announcements",
    ]),
    ("Content Creator", "cc@gmail.com", "cc@123", AMBER, RGBColor(0xFF,0xFB,0xEB), RGBColor(0xFC,0xD3,0x4D), [
        "Upload and manage learning materials",
        "Edit or delete own content items",
        "Manage the content library",
        "View performance analytics for content",
        "Access About and Contact pages",
    ]),
    ("Student", "st@gmail.com", "st@123", GREEN, RGBColor(0xEC,0xFD,0xF5), RGBColor(0x6E,0xE7,0xB7), [
        "Browse all published courses",
        "Enroll in courses of interest",
        "Access enrolled courses via My Learning tab",
        "View and submit assigned work",
        "Track personal learning progress",
        "Read platform announcements",
    ]),
]

positions = [(0.3, 1.25), (3.55, 1.25), (6.8, 1.25), (10.05, 1.25)]

for i, (role, email, pwd, color, bg_c, border_c, bullets) in enumerate(roles):
    lx, ty = positions[i]
    box(s3, lx, ty, 3.0, 6.05, bg_c, border_c, 1.5)
    # Role title bar
    tb = box(s3, lx, ty, 3.0, 0.5, color, color, 0)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = role
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = WHITE

    txt(s3, email, lx+0.1, ty+0.55, 2.8, 0.25, size=10,
        color=DKGREY, italic=True)
    txt(s3, "Pass: " + pwd, lx+0.1, ty+0.78, 2.8, 0.25, size=9,
        color=DKGREY, italic=True)

    cur = ty + 1.05
    for b in bullets:
        bbox = s3.shapes.add_textbox(Inches(lx+0.15), Inches(cur), Inches(2.75), Inches(0.3))
        bbox.word_wrap = True
        btf = bbox.text_frame
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        br = bp.add_run()
        br.text = "- " + b
        br.font.size = Pt(10)
        br.font.color.rgb = BLACK
        cur += 0.32


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — APPLICATION FLOW
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s4 = prs.slides.add_slide(BLANK)
bg(s4, GREY)
header_bar(s4, "Application Flow", "How a user moves from first visit to their dashboard")

# Main flow row
nodes_main = ["User Visits Site", "Auth Page", "Login / Sign Up", "Credentials Verified", "Role Dashboard"]
start_x = 0.5
node_w = 2.0
gap = 0.45
y_row = 1.8

for i, label in enumerate(nodes_main):
    lx = start_x + i*(node_w + gap)
    fc = NAVY if i in [0,4] else LBLUE
    bc = NAVY if i in [0,4] else BLUE
    tc = WHITE if i in [0,4] else NAVY
    flow_box(s4, label, lx, y_row, node_w, 0.65, fc, bc, tc)
    if i < len(nodes_main)-1:
        arrow(s4, lx+node_w+0.05, y_row+0.31, gap-0.1)

# Branch label
txt(s4, "Role-based redirection after successful authentication:", 0.5, 2.75, 12, 0.35,
    size=13, bold=True, color=NAVY)

# Four role branches
branch_data = [
    ("Admin",           "Full platform control\nUsers, Courses, Analytics, Settings",    RED,    RGBColor(0xFE,0xF2,0xF2), RGBColor(0xFC,0xA5,0xA5)),
    ("Instructor",      "Teaching tools\nCourses, Assignments, Analytics",              PURPLE, RGBColor(0xF5,0xF3,0xFF), RGBColor(0xC4,0xB5,0xFD)),
    ("Content Creator", "Media management\nContent Library, Performance",               AMBER,  RGBColor(0xFF,0xFB,0xEB), RGBColor(0xFC,0xD3,0x4D)),
    ("Student",         "Learning journey\nBrowse, Enroll, Submit, Progress",           GREEN,  RGBColor(0xEC,0xFD,0xF5), RGBColor(0x6E,0xE7,0xB7)),
]

bw = 2.9
for i, (role, desc, color, bgc, bdc) in enumerate(branch_data):
    lx = 0.5 + i*(bw+0.3)
    b = box(s4, lx, 3.2, bw, 1.3, bgc, bdc, 1.5)
    tb = box(s4, lx, 3.2, bw, 0.45, color, color, 0)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = role
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WHITE
    txt(s4, desc, lx+0.1, 3.68, bw-0.2, 0.85, size=11, color=BLACK)

# Sub-flows
txt(s4, "Authentication Sub-Flows", 0.5, 4.7, 12, 0.35, size=13, bold=True, color=NAVY)

sub_flows = [
    ("Sign Up Flow", ["Select role (Admin/Instructor/Creator/Student)", "Fill name, email, password, confirm password", "Enter staff code if non-student role", "Account created and auto-logged in to dashboard"]),
    ("Login Flow", ["Enter registered email address", "Enter account password", "System validates credentials", "Redirected to role-specific dashboard"]),
    ("Logout Flow", ["Click Sign Out button in sidebar footer", "Current session is cleared from storage", "User redirected to the authentication page"]),
]

for i, (title, steps) in enumerate(sub_flows):
    lx = 0.5 + i*4.28
    box(s4, lx, 5.1, 4.0, 2.3, WHITE, RGBColor(0xDD,0xE6,0xF0), 1)
    txt(s4, title, lx+0.15, 5.18, 3.7, 0.32, size=12, bold=True, color=NAVY)
    for j, step in enumerate(steps):
        txt(s4, f"{j+1}.  {step}", lx+0.15, 5.54+j*0.45, 3.7, 0.4, size=10, color=BLACK)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — UI SCREENS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s5 = prs.slides.add_slide(BLANK)
bg(s5, GREY)
header_bar(s5, "UI Screens", "All major screens accessible in the application")

screen_data = [
    ("Login Page",           "All Roles",       BLUE,   "Sliding panel layout with login and signup forms, password strength indicator, show/hide toggle"),
    ("Signup - Role Picker", "All Roles",       BLUE,   "Step 1: user selects their role from four options. Step 2: fills name, email, password, and staff code if required"),
    ("Admin Dashboard",      "Admin",           RED,    "Overview statistics: total users, courses, enrollments, assignments. Quick-access cards for core modules"),
    ("User Management",      "Admin",           RED,    "Full table of all users with add, edit, deactivate, and delete operations. Filterable and searchable"),
    ("Courses Page",         "All Roles",       BLUE,   "Admin: create and manage courses. Student: browse and enroll. Instructor: view own courses. Filter by category"),
    ("Assignments",          "Admin, Student",  PURPLE, "Admin creates assignments with due dates and max scores. Students view and submit. Admin grades submissions"),
    ("Announcements",        "All Roles",       GREEN,  "Platform-wide communication board. Admin creates posts visible to all roles on the platform"),
    ("Analytics",            "Admin",           AMBER,  "Charts and statistics for course enrollments, user distribution, submissions, and platform activity"),
    ("Settings",             "Admin",           BLACK,  "Edit platform name, contact email, phone number, and the About page content directly from this screen"),
    ("Student Dashboard",    "Student",         GREEN,  "Shows enrolled courses, upcoming assignments, learning progress, and announcements in one view"),
    ("Browse Courses",       "Student",         BLUE,   "Grid of all published courses with category filters, enrollment status badges, and course detail modal"),
    ("My Learning",          "Student",         GREEN,  "Filtered view showing only courses the student is enrolled in, with quick access to continue learning"),
]

cols = 4
rows = 3
card_w = 3.0
card_h = 1.75
gap_x = 0.2
gap_y = 0.2
start_x = 0.3
start_y = 1.25

for idx, (name, role, color, desc) in enumerate(screen_data):
    col = idx % cols
    row = idx // cols
    lx = start_x + col*(card_w+gap_x)
    ty = start_y + row*(card_h+gap_y)

    # Card bg
    box(s5, lx, ty, card_w, card_h, WHITE, RGBColor(0xDD,0xE6,0xF0), 1)
    # Color top strip
    strip = box(s5, lx, ty, card_w, 0.38, color, color, 0)
    tf = strip.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = name
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # Role badge
    txt(s5, "Role: " + role, lx+0.1, ty+0.42, card_w-0.2, 0.25, size=9, color=DKGREY, bold=True)
    # Description
    txt(s5, desc, lx+0.1, ty+0.65, card_w-0.2, 1.05, size=9, color=BLACK)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 — NAVIGATION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s6 = prs.slides.add_slide(BLANK)
bg(s6, GREY)
header_bar(s6, "Navigation", "Role-based sidebar navigation — each role sees only their permitted pages")

nav_data = [
    ("Admin", RED, [
        "Dashboard - Platform overview and statistics",
        "User Management - Add, edit, delete users",
        "All Courses - Create and manage all courses",
        "Assignments - Create and grade assignments",
        "Announcements - Post platform messages",
        "Analytics - View platform-wide data",
        "Platform Settings - Configure the system",
        "About Us - Editable about page",
        "Contact Us - Contact information page",
    ]),
    ("Instructor", PURPLE, [
        "Dashboard - Teaching overview",
        "My Courses - View own created courses",
        "Assignments - Manage course assignments",
        "Announcements - Read platform messages",
        "Analytics - View course performance data",
        "About Us - About page",
        "Contact Us - Contact page",
    ]),
    ("Content Creator", AMBER, [
        "Dashboard - Content overview",
        "Content Library - Upload and manage materials",
        "Performance - View content analytics",
        "About Us - About page",
        "Contact Us - Contact page",
    ]),
    ("Student", GREEN, [
        "Dashboard - Learning overview",
        "Browse Courses - Discover all published courses",
        "My Learning - Access enrolled courses only",
        "Assignments - View and submit work",
        "Announcements - Read platform messages",
        "My Progress - Track personal analytics",
        "About Us - About page",
        "Contact Us - Contact page",
    ]),
]

col_w = 3.0
gap_c = 0.27
sx = 0.3

for i, (role, color, navs) in enumerate(nav_data):
    lx = sx + i*(col_w+gap_c)
    box(s6, lx, 1.25, col_w, 6.1, WHITE, RGBColor(0xDD,0xE6,0xF0), 1)
    strip = box(s6, lx, 1.25, col_w, 0.48, color, color, 0)
    tf = strip.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = role
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WHITE

    cur_y = 1.82
    for nav in navs:
        parts = nav.split(" - ", 1)
        nb = s6.shapes.add_textbox(Inches(lx+0.15), Inches(cur_y), Inches(col_w-0.25), Inches(0.38))
        nb.word_wrap = True
        ntf = nb.text_frame
        ntf.word_wrap = True
        np_ = ntf.paragraphs[0]
        nr1 = np_.add_run()
        nr1.text = parts[0]
        nr1.font.size = Pt(10)
        nr1.font.bold = True
        nr1.font.color.rgb = color
        if len(parts) > 1:
            nr2 = np_.add_run()
            nr2.text = "  " + parts[1]
            nr2.font.size = Pt(10)
            nr2.font.bold = False
            nr2.font.color.rgb = BLACK
        cur_y += 0.42

# Mobile note
txt(s6, "On mobile devices, the sidebar is hidden by default. A hamburger button (top-left) reveals it as a slide-in drawer with a dark overlay. Tapping any menu item navigates and closes the sidebar automatically.", 0.3, 7.1, 12.73, 0.38, size=10, color=DKGREY, italic=True)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 — FUNCTIONAL WORKFLOW
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s7 = prs.slides.add_slide(BLANK)
bg(s7, GREY)
header_bar(s7, "Functional Workflow", "Step-by-step operations for core platform functions")

workflows = [
    ("User Registration", NAVY, [
        "Visit the application and click Create one on the auth page",
        "Select a role: Admin, Instructor, Content Creator, or Student",
        "Fill in full name, email address, password, and confirm password",
        "Enter staff code (required for Instructor, Content Creator, Admin roles)",
        "Click Create Account to complete registration",
        "Account is created and user is automatically logged into their dashboard",
    ]),
    ("Admin: Create and Publish Course", RED, [
        "Log in with admin credentials",
        "Navigate to All Courses in the sidebar",
        "Click the New Course button in the top-right corner",
        "Enter title, category, level, duration, number of lessons, and description",
        "Set the status to Published so students can see the course",
        "Click Create Course - the course is now visible to all enrolled students",
    ]),
    ("Student: Browse and Enroll", GREEN, [
        "Log in with student credentials",
        "Navigate to Browse Courses in the sidebar",
        "Browse the grid of published courses, filter by category if needed",
        "Click any course card to open the course detail modal",
        "Review course information including lessons, duration, and description",
        "Click Enroll Now - the course appears in My Learning immediately",
    ]),
    ("Assignment: Full Lifecycle", PURPLE, [
        "Admin creates an assignment linked to a specific course with a due date",
        "Students enrolled in that course see the assignment in their Assignments page",
        "Student writes an answer in the text field and clicks Submit",
        "Admin and Instructor receive a notification about the new submission",
        "Admin opens the submission and enters a score and written feedback",
        "Student can view their grade and feedback in the Assignments section",
    ]),
    ("User Management (Admin)", AMBER, [
        "Admin navigates to User Management in the sidebar",
        "View a table of all registered users with name, email, role, and status",
        "Click Add User to create a new account with a specific role and password",
        "Click the edit icon to modify name, role, or account status of any user",
        "Toggle a user to inactive to suspend access without deleting their data",
        "Click the delete icon and confirm to permanently remove a user account",
    ]),
    ("Content Upload (Creator)", RGBColor(0x0E,0x74,0x90), [
        "Log in with Content Creator credentials",
        "Navigate to Content Library in the sidebar",
        "Click Add Content to open the upload form",
        "Enter title, content type (video, document, PDF), URL, and any notes",
        "Associate the material with a specific course or leave it as general",
        "Click Save - the content appears in the library and is viewable by others",
    ]),
]

wf_w = 4.0
wf_gap = 0.27
wf_sx = 0.3
for i, (title, color, steps) in enumerate(workflows):
    col = i % 3
    row = i // 3
    lx = wf_sx + col*(wf_w+wf_gap)
    ty = 1.25 + row*3.1
    box(s7, lx, ty, wf_w, 3.0, WHITE, RGBColor(0xDD,0xE6,0xF0), 1)
    strip = box(s7, lx, ty, wf_w, 0.44, color, color, 0)
    tf = strip.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE
    for j, step in enumerate(steps):
        txt(s7, f"{j+1}.  {step}", lx+0.15, ty+0.5+j*0.41, wf_w-0.25, 0.38, size=9.5, color=BLACK)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 — LINKS AND TEAM
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s8 = prs.slides.add_slide(BLANK)
bg(s8, NAVY)

txt(s8, "Digital Black Board — LMS", 0.5, 0.6, 12.33, 0.7, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s8, "TEAM 18  |  KL University", 0.5, 1.35, 12.33, 0.4, size=14, color=RGBColor(0xA8,0xC4,0xE0), align=PP_ALIGN.CENTER)

# Link boxes
links = [
    ("Live Application", "https://digital-blackboard.netlify.app", BLUE),
    ("Source Code (GitHub)", "https://github.com/NALIN-9/LMS-PROJECT", RGBColor(0x16,0x18,0x22)),
]
for i, (label, url, color) in enumerate(links):
    lx = 1.5 + i*5.5
    box(s8, lx, 2.0, 5.0, 0.9, color, color, 0)
    txt(s8, label, lx+0.2, 2.08, 4.6, 0.3, size=12, bold=True, color=WHITE)
    txt(s8, url, lx+0.2, 2.42, 4.6, 0.45, size=11, color=RGBColor(0xBF,0xDB,0xFE))

# Credentials
txt(s8, "Login Credentials", 0.5, 3.1, 12.33, 0.35, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
creds = [
    ("Admin",           "admin@gmail.com", "Admin@123", RED),
    ("Instructor",      "ins@gmail.com",   "ins@123",   PURPLE),
    ("Content Creator", "cc@gmail.com",    "cc@123",    AMBER),
    ("Student",         "st@gmail.com",    "st@123",    GREEN),
]
for i, (role, email, pwd, color) in enumerate(creds):
    lx = 0.4 + i*3.2
    box(s8, lx, 3.55, 3.0, 1.25, RGBColor(0x14,0x1E,0x33), color, 1.5)
    txt(s8, role, lx+0.15, 3.62, 2.7, 0.3, size=12, bold=True, color=color)
    txt(s8, email, lx+0.15, 3.96, 2.7, 0.25, size=10, color=WHITE)
    txt(s8, f"Password: {pwd}", lx+0.15, 4.2, 2.7, 0.25, size=10, color=RGBColor(0xA8,0xC4,0xE0))

txt(s8, "Staff Code for signup (non-student roles):  DBBLMS", 0.5, 4.95, 12.33, 0.35,
    size=12, color=RGBColor(0xBF,0xDB,0xFE), align=PP_ALIGN.CENTER)

# Team
box(s8, 0.5, 5.45, 12.33, 1.75, RGBColor(0x14,0x1E,0x33), BLUE, 1)
txt(s8, "Team Members", 0.7, 5.55, 11.9, 0.35, size=13, bold=True, color=WHITE)
members = [
    "KANTAMANI NALIN KUMAR      (2400030332)",
    "GHANTA NAGA PRASANTH BABU  (2400030196)",
    "KARANAM BHARGAV            (2400031923)",
]
for i, m in enumerate(members):
    txt(s8, m, 0.7, 5.95+i*0.32, 7.0, 0.3, size=11, color=RGBColor(0xBF,0xDB,0xFE))

txt(s8, "Guide: Dr. NAGARJUNA KARYEMSETTY  |  KL University", 7.5, 5.95, 5.0, 0.65,
    size=11, color=RGBColor(0xA8,0xC4,0xE0))


# ── SAVE ──────────────────────────────────────────────────────
out = r"C:\Users\nalin\Downloads\DBB_Prototype_Presentation.pptx"
prs.save(out)
print("Saved:", out)
